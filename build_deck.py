"""
Build a 1-hour Prompt Engineering presentation.
Output: prompt_engineering_workshop.pptx
References:
  - https://www.promptingguide.ai/
  - https://www.promptingguide.ai/introduction/tips
  - https://www.promptingguide.ai/techniques
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ---------- Theme ----------
NAVY      = RGBColor(0x0B, 0x1F, 0x3A)
INK       = RGBColor(0x10, 0x18, 0x28)
BG        = RGBColor(0xF7, 0xF8, 0xFA)
ACCENT    = RGBColor(0xFF, 0x9F, 0x1C)   # warm orange
ACCENT2   = RGBColor(0x2E, 0xC4, 0xB6)   # teal
MUTED     = RGBColor(0x55, 0x6B, 0x83)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LINE      = RGBColor(0xE5, 0xE7, 0xEB)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=INK,
                bold_lead=False, line_spacing=1.15):
    """items: list of strings; supports 'Lead: rest' to bold the lead."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        # Bullet glyph
        b = p.add_run()
        b.text = "•  "
        b.font.size = Pt(size)
        b.font.color.rgb = ACCENT
        b.font.bold = True
        b.font.name = "Calibri"

        if bold_lead and ":" in item:
            lead, rest = item.split(":", 1)
            r1 = p.add_run()
            r1.text = lead + ":"
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r1.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = "Calibri"
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = "Calibri"
    return tb


def base_slide(title, eyebrow=None):
    s = prs.slides.add_slide(BLANK)
    # Background
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, BG)
    # Top accent bar
    add_rect(s, 0, 0, prs.slide_width, Inches(0.18), ACCENT)
    # Eyebrow
    if eyebrow:
        add_text(s, Inches(0.6), Inches(0.35), Inches(8), Inches(0.35),
                 eyebrow.upper(), size=11, bold=True, color=ACCENT2)
    # Title
    add_text(s, Inches(0.6), Inches(0.65), Inches(12.1), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    # Underline accent
    add_rect(s, Inches(0.6), Inches(1.35), Inches(0.6), Inches(0.05), ACCENT)
    # Footer
    add_text(s, Inches(0.6), Inches(7.15), Inches(8), Inches(0.3),
             "Prompt Engineering Workshop  ·  60-minute session",
             size=9, color=MUTED)
    add_text(s, Inches(11.1), Inches(7.15), Inches(1.7), Inches(0.3),
             f"Slide {len(prs.slides)}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    return s


def code_box(slide, x, y, w, h, text, *, size=12):
    add_rect(slide, x, y, w, h, RGBColor(0x0E, 0x1B, 0x2C), line=LINE)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.splitlines()):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        r.font.color.rgb = RGBColor(0xE6, 0xED, 0xF3)
    return tb


def card(slide, x, y, w, h, title, body_lines, *, accent=ACCENT2):
    add_rect(slide, x, y, w, h, WHITE, line=LINE)
    add_rect(slide, x, y, Inches(0.08), h, accent)
    add_text(slide, x + Inches(0.25), y + Inches(0.15), w - Inches(0.4), Inches(0.5),
             title, size=16, bold=True, color=NAVY)
    add_bullets(slide, x + Inches(0.25), y + Inches(0.65), w - Inches(0.4),
                h - Inches(0.7), body_lines, size=12, color=INK)


# =====================================================================
# 1. Title
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
# Decorative shapes
add_rect(s, Inches(10.5), 0, Inches(2.83), Inches(7.5), ACCENT2)
add_rect(s, Inches(10.5), Inches(4), Inches(2.83), Inches(3.5), ACCENT)

add_text(s, Inches(0.7), Inches(1.0), Inches(9), Inches(0.4),
         "WORKSHOP · 60 MINUTES", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(1.5), Inches(10), Inches(1.4),
         "Prompt Engineering", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(2.7), Inches(10), Inches(1.0),
         "Core Concepts & Best Practices", size=30, color=WHITE)
add_text(s, Inches(0.7), Inches(4.3), Inches(9), Inches(0.5),
         "Designing prompts that are clear, reliable, and production-ready",
         size=18, color=RGBColor(0xCB, 0xD5, 0xE1))

# Reference strip
add_rect(s, Inches(0.7), Inches(6.3), Inches(9.4), Inches(0.7),
         RGBColor(0x12, 0x2C, 0x4E))
add_text(s, Inches(0.9), Inches(6.42), Inches(9), Inches(0.5),
         "Reference: promptingguide.ai · DAIR.AI",
         size=14, color=WHITE)


# =====================================================================
# 2. Agenda
# =====================================================================
s = base_slide("Agenda", eyebrow="What we'll cover")

agenda = [
    ("01", "Foundations",          "What prompting is, why it matters, mental model"),
    ("02", "Anatomy & Settings",   "Prompt elements, temperature, top-p, max tokens"),
    ("03", "Core Techniques",      "Zero-shot, few-shot, chain-of-thought, self-consistency"),
    ("04", "Advanced Patterns",    "Prompt chaining, ToT, ReAct, RAG, agents"),
    ("05", "Risks & Safety",       "Hallucination, bias, prompt injection, defenses"),
    ("06", "Workflow & Evaluation","Iteration, testing, production discipline"),
]

start_y = 1.75
for i, (num, title, desc) in enumerate(agenda):
    y = Inches(start_y + i * 0.85)
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.75), WHITE, line=LINE)
    add_rect(s, Inches(0.6), y, Inches(0.08), Inches(0.75), ACCENT)
    add_text(s, Inches(0.85), y + Inches(0.12), Inches(0.9), Inches(0.5),
             num, size=20, bold=True, color=ACCENT2)
    add_text(s, Inches(1.85), y + Inches(0.07), Inches(3.3), Inches(0.4),
             title, size=17, bold=True, color=NAVY)
    add_text(s, Inches(1.85), y + Inches(0.4), Inches(10.5), Inches(0.35),
             desc, size=13, color=MUTED)


# =====================================================================
# 3. What is prompt engineering?
# =====================================================================
s = base_slide("What is prompt engineering?", eyebrow="Foundations")

add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.9),
         "The discipline of designing inputs that reliably steer language models "
         "toward useful, accurate outputs.",
         size=22, color=INK)

# Three pillars
pillars = [
    ("Clarity",    "Express intent precisely. Remove ambiguity. State the task, the audience, and the format."),
    ("Context",    "Give the model what it needs and nothing more. Include background, examples, and constraints."),
    ("Control",    "Use structure, settings, and patterns to make outputs predictable and verifiable."),
]
for i, (title, body) in enumerate(pillars):
    x = Inches(0.6 + i * 4.15)
    y = Inches(3.2)
    add_rect(s, x, y, Inches(3.95), Inches(2.6), WHITE, line=LINE)
    add_rect(s, x, y, Inches(3.95), Inches(0.55), ACCENT2)
    add_text(s, x + Inches(0.25), y + Inches(0.08), Inches(3.7), Inches(0.45),
             title, size=18, bold=True, color=WHITE)
    add_text(s, x + Inches(0.25), y + Inches(0.75), Inches(3.5), Inches(1.7),
             body, size=14, color=INK)

add_text(s, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.6),
         "Good prompts are written like good specs: specific, testable, and iterated.",
         size=15, bold=True, color=MUTED)


# =====================================================================
# 4. Why it matters
# =====================================================================
s = base_slide("Why it matters", eyebrow="Foundations")

stats = [
    ("40-60%",  "Output accuracy lift from structured vs. ad-hoc prompts*"),
    ("10x",     "Faster iteration than fine-tuning for most use cases"),
    ("0 GPUs",  "No training infrastructure required to ship improvements"),
]
for i, (n, desc) in enumerate(stats):
    x = Inches(0.6 + i * 4.15)
    add_rect(s, x, Inches(1.8), Inches(3.95), Inches(2.0), WHITE, line=LINE)
    add_text(s, x, Inches(1.95), Inches(3.95), Inches(0.9),
             n, size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), Inches(2.85), Inches(3.35), Inches(0.9),
             desc, size=13, color=INK, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(0.5),
         "What prompt engineering buys you", size=18, bold=True, color=NAVY)

card(s, Inches(0.6),  Inches(4.65), Inches(3.95), Inches(2.2),
     "Reliability",
     ["Reduce hallucinations and off-topic answers",
      "Stable outputs across runs",
      "Predictable formats for downstream code"])
card(s, Inches(4.7),  Inches(4.65), Inches(3.95), Inches(2.2),
     "Speed",
     ["Ship behavior changes in minutes",
      "Lower cost than fine-tuning",
      "Faster A/B between models"])
card(s, Inches(8.8),  Inches(4.65), Inches(3.95), Inches(2.2),
     "Capability",
     ["Unlock reasoning and tool use",
      "Compose simple steps into workflows",
      "Adapt one model to many tasks"])

add_text(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.3),
         "* Stanford HAI, 2025 (industry-reported range)", size=9, color=MUTED)


# =====================================================================
# 5. Mental model: how LLMs respond
# =====================================================================
s = base_slide("Mental model: how LLMs respond", eyebrow="Foundations")

add_text(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.6),
         "LLMs are next-token predictors conditioned on everything in the context window.",
         size=18, bold=True, color=INK)

# Flow diagram
y = Inches(2.8)
boxes = [
    ("Context",   "System + user\nmessages, examples,\ntools, retrieved docs"),
    ("Model",     "Pretrained weights\n+ tokenizer\n+ sampling settings"),
    ("Tokens",    "One-by-one\nprobabilistic\ngeneration"),
    ("Output",    "Text, code,\nstructured data,\nor tool calls"),
]
bw = Inches(2.7); bh = Inches(2.0); gap = Inches(0.35)
total_w = bw * 4 + gap * 3
start_x = (prs.slide_width - total_w) // 2
for i, (t, body) in enumerate(boxes):
    x = start_x + (bw + gap) * i
    add_rect(s, x, y, bw, bh, WHITE, line=LINE)
    add_rect(s, x, y, bw, Inches(0.5), NAVY)
    add_text(s, x, y + Inches(0.05), bw, Inches(0.4),
             t, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(0.65), bw - Inches(0.4),
             bh - Inches(0.7), body, size=13, color=INK, align=PP_ALIGN.CENTER)
    if i < 3:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   x + bw + Inches(0.02),
                                   y + bh / 2 - Inches(0.18),
                                   Inches(0.31), Inches(0.36))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

add_text(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.5),
         "Implications for how you prompt", size=16, bold=True, color=NAVY)

bullets = [
    "Everything in the context shapes output: ordering, examples, formatting, and stray instructions all matter.",
    "Models do not know what they do not see. Provide facts; do not assume the model has them.",
    "Sampling settings change behavior. Pin them when you measure prompt changes.",
    "Long contexts dilute attention. Keep prompts focused; trim what you don't need.",
]
add_bullets(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(1.4),
            bullets, size=13)


# =====================================================================
# 6. Anatomy of a prompt
# =====================================================================
s = base_slide("Anatomy of a prompt", eyebrow="Anatomy & Settings")

add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.5),
         "A prompt typically contains four elements. You don't always need all four.",
         size=15, color=MUTED)

elements = [
    ("Instruction",     "The task: 'classify', 'summarize', 'extract', 'translate'."),
    ("Context",         "Background, persona, constraints, retrieved docs."),
    ("Input data",      "The specific input the model should act on."),
    ("Output indicator","The expected format, schema, or example shape."),
]
for i, (t, body) in enumerate(elements):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 4.0)
    y = Inches(2.4 + row * 1.45)
    add_rect(s, x, y, Inches(3.85), Inches(1.3), WHITE, line=LINE)
    add_text(s, x + Inches(0.2), y + Inches(0.12), Inches(3.6), Inches(0.4),
             t, size=15, bold=True, color=ACCENT2)
    add_text(s, x + Inches(0.2), y + Inches(0.55), Inches(3.6), Inches(0.7),
             body, size=12, color=INK)

# Example code box
code_box(s, Inches(8.7), Inches(2.4), Inches(4.0), Inches(4.4),
"""### Instruction ###
Classify the sentiment of the
review as positive, neutral,
or negative.

### Context ###
You are a support analyst.
Be concise.

### Input ###
Review: "Shipping took forever
but the product is excellent."

### Output ###
Sentiment: <one word>
Reason: <one sentence>""", size=12)


# =====================================================================
# 7. LLM Settings
# =====================================================================
s = base_slide("LLM settings that matter", eyebrow="Anatomy & Settings")

settings = [
    ("Temperature",
     "Controls randomness.\nLower = deterministic.\nHigher = creative.",
     "0.0 - 0.3 for extraction\n0.7+ for ideation"),
    ("Top-p (nucleus)",
     "Caps the probability mass\nsampled from. Use one of\ntemperature OR top-p.",
     "0.9 is a common default"),
    ("Max tokens",
     "Hard limit on output length.\nProtects cost and latency.",
     "Right-size per task"),
    ("Stop sequences",
     "Force generation to halt at\nspecific markers.",
     "Useful for structured output"),
    ("System prompt",
     "Persistent instructions and\npersona for the conversation.",
     "Set once, reuse"),
    ("Seed (when supported)",
     "Reproducible sampling for\ntesting and evaluation.",
     "Pin during prompt iteration"),
]
for i, (t, body, hint) in enumerate(settings):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.15)
    y = Inches(1.85 + row * 2.6)
    add_rect(s, x, y, Inches(3.95), Inches(2.4), WHITE, line=LINE)
    add_rect(s, x, y, Inches(3.95), Inches(0.5), NAVY)
    add_text(s, x + Inches(0.2), y + Inches(0.08), Inches(3.6), Inches(0.4),
             t, size=15, bold=True, color=WHITE)
    add_text(s, x + Inches(0.2), y + Inches(0.65), Inches(3.6), Inches(1.1),
             body, size=12, color=INK)
    add_text(s, x + Inches(0.2), y + Inches(1.85), Inches(3.6), Inches(0.5),
             hint, size=11, bold=True, color=ACCENT2)


# =====================================================================
# 8. General tips for designing prompts
# =====================================================================
s = base_slide("General tips for designing prompts", eyebrow="Best Practices")

tips = [
    ("Start simple, iterate",
     "Begin with a minimal prompt. Add elements only when they fix a real failure mode."),
    ("Be specific",
     "State the task, audience, format, and constraints. Vague prompts produce vague outputs."),
    ("Show, don't only tell",
     "Examples beat adjectives. A few well-chosen demonstrations anchor format and style."),
    ("Say what to do, not what to avoid",
     "'Recommend from the top trending list' beats 'don't ask for preferences'."),
    ("Use clear separators",
     "Delimit sections with markers like ###, XML tags, or fenced blocks."),
    ("Decompose hard tasks",
     "Split big jobs into smaller prompts you can test and chain together."),
]
for i, (t, body) in enumerate(tips):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.15)
    y = Inches(1.85 + row * 1.7)
    add_rect(s, x, y, Inches(5.95), Inches(1.55), WHITE, line=LINE)
    add_rect(s, x, y, Inches(0.08), Inches(1.55), ACCENT)
    add_text(s, x + Inches(0.25), y + Inches(0.15), Inches(5.6), Inches(0.45),
             t, size=15, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.65), Inches(5.6), Inches(0.85),
             body, size=12, color=INK)


# =====================================================================
# 9. Zero-shot prompting
# =====================================================================
s = base_slide("Zero-shot prompting", eyebrow="Core Techniques")

add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.6),
         "Ask the model to perform a task with no examples. "
         "Works well when the task is common and the instruction is precise.",
         size=15, color=INK)

add_bullets(s, Inches(0.6), Inches(2.55), Inches(6.2), Inches(2.6),
            ["Best for: classification, summarization, simple extraction, translation",
             "Pros: Fast, cheap, no example curation",
             "Cons: Brittle for unusual formats or domain-specific style",
             "Tip: Lead with a strong verb and a clear output schema"],
            size=14)

code_box(s, Inches(7.0), Inches(2.55), Inches(5.7), Inches(2.5),
"""Classify the text into neutral,
negative, or positive.

Text: I think the vacation is okay.

Sentiment:""")

add_text(s, Inches(0.6), Inches(5.55), Inches(12.1), Inches(0.5),
         "When zero-shot fails, the next move is usually few-shot.",
         size=14, bold=True, color=ACCENT2)


# =====================================================================
# 10. Few-shot prompting
# =====================================================================
s = base_slide("Few-shot prompting", eyebrow="Core Techniques")

add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.6),
         "Show the model a few input/output pairs so it can pattern-match the task.",
         size=15, color=INK)

add_bullets(s, Inches(0.6), Inches(2.55), Inches(6.2), Inches(3.4),
            ["Use 2-5 examples that cover the variation you expect",
             "Match the format you want exactly. The model copies what it sees.",
             "Diversity > quantity. Avoid near-duplicates.",
             "Order matters. Recent examples tend to weigh more.",
             "Watch out: bad examples teach bad behavior."],
            size=14)

code_box(s, Inches(7.0), Inches(2.55), Inches(5.7), Inches(3.6),
"""Classify sentiment.

Text: The food was great.
Sentiment: positive

Text: Service was slow but ok.
Sentiment: neutral

Text: I will not return.
Sentiment: negative

Text: Loved every minute.
Sentiment:""")


# =====================================================================
# 11. Chain-of-thought (CoT)
# =====================================================================
s = base_slide("Chain-of-thought prompting", eyebrow="Core Techniques")

add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.6),
         "Encourage the model to reason step by step before answering. "
         "Improves accuracy on math, logic, and multi-hop questions.",
         size=15, color=INK)

add_bullets(s, Inches(0.6), Inches(2.55), Inches(6.2), Inches(3.6),
            ["Zero-shot CoT: append \"Let's think step by step.\"",
             "Few-shot CoT: include worked examples that show the reasoning",
             "For reasoning models, keep prompts simpler — they already think",
             "Trade-off: more tokens, higher latency, higher cost",
             "Pair with self-consistency for best results"],
            size=14)

code_box(s, Inches(7.0), Inches(2.55), Inches(5.7), Inches(3.6),
"""Q: Roger has 5 tennis balls. He buys
2 cans, each with 3 balls. How many
balls does he have now?

A: Let's think step by step.
Roger starts with 5.
2 cans x 3 balls = 6.
5 + 6 = 11.

Answer: 11""")


# =====================================================================
# 12. Self-consistency
# =====================================================================
s = base_slide("Self-consistency", eyebrow="Core Techniques")

add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.7),
         "Sample multiple reasoning paths at higher temperature, then take a "
         "majority vote on the final answer.",
         size=15, color=INK)

# Diagram
y = Inches(2.9)
add_rect(s, Inches(0.8), y, Inches(2.4), Inches(1.0), WHITE, line=LINE)
add_text(s, Inches(0.8), y + Inches(0.3), Inches(2.4), Inches(0.5),
         "Prompt + question", size=14, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)

paths = ["Path 1 → 11", "Path 2 → 11", "Path 3 → 9", "Path 4 → 11"]
for i, p in enumerate(paths):
    x = Inches(4.2)
    yy = Inches(2.0 + i * 0.85)
    add_rect(s, x, yy, Inches(3.2), Inches(0.6), WHITE, line=LINE)
    add_text(s, x, yy + Inches(0.13), Inches(3.2), Inches(0.4),
             p, size=14, color=INK, align=PP_ALIGN.CENTER)

add_rect(s, Inches(8.9), y, Inches(3.6), Inches(1.0), ACCENT2)
add_text(s, Inches(8.9), y + Inches(0.25), Inches(3.6), Inches(0.5),
         "Majority vote → 11", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)

add_bullets(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.6),
            ["Robust against reasoning slips on hard problems",
             "Cost scales linearly with sample count — typically 5-20 samples",
             "Works best when the final answer is short and verifiable",
             "Pair with CoT to generate diverse reasoning traces"],
            size=14)


# =====================================================================
# 13. Generate knowledge & prompt chaining
# =====================================================================
s = base_slide("Generate knowledge & prompt chaining", eyebrow="Advanced Patterns")

# Two-column layout
add_rect(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(5.0), WHITE, line=LINE)
add_rect(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(0.55), NAVY)
add_text(s, Inches(0.8), Inches(1.88), Inches(5.6), Inches(0.4),
         "Generate Knowledge", size=16, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.5), Inches(5.6), Inches(0.7),
         "Ask the model to first produce relevant facts, then answer using them.",
         size=13, color=INK)
add_bullets(s, Inches(0.8), Inches(3.3), Inches(5.6), Inches(3.3),
            ["Step 1: 'List 3 facts about X.'",
             "Step 2: 'Using those facts, answer Y.'",
             "Reduces shallow guessing on knowledge tasks",
             "Works even better when facts come from retrieval"],
            size=13)

add_rect(s, Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.0), WHITE, line=LINE)
add_rect(s, Inches(6.8), Inches(1.8), Inches(6.0), Inches(0.55), NAVY)
add_text(s, Inches(7.0), Inches(1.88), Inches(5.6), Inches(0.4),
         "Prompt Chaining", size=16, bold=True, color=WHITE)
add_text(s, Inches(7.0), Inches(2.5), Inches(5.6), Inches(0.7),
         "Break a hard task into a sequence of small, focused prompts.",
         size=13, color=INK)
add_bullets(s, Inches(7.0), Inches(3.3), Inches(5.6), Inches(3.3),
            ["Each step has a clear input and output contract",
             "Easier to test, debug, and swap models per step",
             "Common pattern: extract → transform → format",
             "Keeps any single context small and on-task"],
            size=13)


# =====================================================================
# 14. Tree of thoughts
# =====================================================================
s = base_slide("Tree of Thoughts (ToT)", eyebrow="Advanced Patterns")

add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.7),
         "Explore multiple intermediate steps as a tree, evaluate each branch, "
         "and search for the best path.",
         size=15, color=INK)

# Simple tree visual
import math
cx = Inches(3.5); cy = Inches(2.9)
node_w = Inches(1.6); node_h = Inches(0.55)

def node(x, y, label, accent=NAVY):
    add_rect(s, x, y, node_w, node_h, WHITE, line=LINE)
    add_rect(s, x, y, node_w, Inches(0.08), accent)
    add_text(s, x, y + Inches(0.15), node_w, Inches(0.4),
             label, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Root
node(cx, cy, "Problem", ACCENT)
# Level 1
node(Inches(1.0), Inches(4.0), "Approach A")
node(Inches(3.5), Inches(4.0), "Approach B", ACCENT2)
node(Inches(6.0), Inches(4.0), "Approach C")
# Level 2 (under B)
node(Inches(2.5), Inches(5.2), "Step B.1")
node(Inches(4.6), Inches(5.2), "Step B.2", ACCENT2)
# Level 3
node(Inches(4.6), Inches(6.3), "Solution", ACCENT)

# Connectors (simple lines)
def connect(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = MUTED
    line.line.width = Pt(1.25)

connect(s, cx + node_w/2, cy + node_h, Inches(1.0) + node_w/2, Inches(4.0))
connect(s, cx + node_w/2, cy + node_h, Inches(3.5) + node_w/2, Inches(4.0))
connect(s, cx + node_w/2, cy + node_h, Inches(6.0) + node_w/2, Inches(4.0))
connect(s, Inches(3.5) + node_w/2, Inches(4.0) + node_h, Inches(2.5) + node_w/2, Inches(5.2))
connect(s, Inches(3.5) + node_w/2, Inches(4.0) + node_h, Inches(4.6) + node_w/2, Inches(5.2))
connect(s, Inches(4.6) + node_w/2, Inches(5.2) + node_h, Inches(4.6) + node_w/2, Inches(6.3))

add_bullets(s, Inches(8.0), Inches(2.9), Inches(4.7), Inches(3.6),
            ["Models propose multiple thoughts per step",
             "Each branch is scored or self-evaluated",
             "Search (BFS/DFS) finds the most promising path",
             "Useful for planning, puzzles, and creative tasks",
             "More expensive than CoT — use when accuracy matters"],
            size=13)


# =====================================================================
# 15. ReAct & tool use
# =====================================================================
s = base_slide("ReAct: reason + act", eyebrow="Advanced Patterns")

add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.6),
         "Interleave reasoning steps with tool calls so the model can fetch "
         "external information and act on it.",
         size=15, color=INK)

# Loop diagram
y = Inches(2.7)
loop = [("Thought", ACCENT2), ("Action", ACCENT), ("Observation", NAVY)]
for i, (t, c) in enumerate(loop):
    x = Inches(0.6 + i * 2.6)
    add_rect(s, x, y, Inches(2.3), Inches(1.1), WHITE, line=LINE)
    add_rect(s, x, y, Inches(2.3), Inches(0.35), c)
    add_text(s, x, y + Inches(0.04), Inches(2.3), Inches(0.3),
             t, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.5), Inches(2.3), Inches(0.55),
             ["Plan next step",
              "Call a tool",
              "Read tool output"][i],
             size=12, color=INK, align=PP_ALIGN.CENTER)
    if i < 2:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   x + Inches(2.32), y + Inches(0.4),
                                   Inches(0.26), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

# Loop arrow
loop_arrow = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW,
                                Inches(0.6), Inches(4.0),
                                Inches(7.2), Inches(0.35))
loop_arrow.fill.solid()
loop_arrow.fill.fore_color.rgb = MUTED
loop_arrow.line.fill.background()
add_text(s, Inches(0.6), Inches(4.0), Inches(7.2), Inches(0.35),
         "repeat until done", size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

code_box(s, Inches(8.5), Inches(2.7), Inches(4.2), Inches(3.6),
"""Thought: I need today's weather.
Action: search("Seattle weather")
Observation: 58°F, light rain
Thought: I have what I need.
Final: Rain expected; bring a
jacket.""", size=11)

add_bullets(s, Inches(0.6), Inches(4.7), Inches(7.5), Inches(2.0),
            ["Foundation pattern for AI agents",
             "Pair with function calling for typed tool I/O",
             "Cap loop iterations to control cost and runaway behavior",
             "Log every step — observability is non-negotiable"],
            size=13)


# =====================================================================
# 16. Retrieval Augmented Generation (RAG)
# =====================================================================
s = base_slide("Retrieval Augmented Generation (RAG)", eyebrow="Advanced Patterns")

add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.6),
         "Inject fresh, domain-specific context into the prompt at query time. "
         "Reduces hallucinations on private or recent information.",
         size=15, color=INK)

# Pipeline
stages = [
    ("Question",  "User query"),
    ("Retrieve",  "Vector or keyword\nsearch over your\nknowledge base"),
    ("Augment",   "Concatenate top\npassages into the\nprompt context"),
    ("Generate",  "Model answers\ngrounded in the\nprovided sources"),
]
y = Inches(2.7); bw = Inches(2.85); bh = Inches(1.7); gap = Inches(0.2)
for i, (t, body) in enumerate(stages):
    x = Inches(0.6) + (bw + gap) * i
    add_rect(s, x, y, bw, bh, WHITE, line=LINE)
    add_rect(s, x, y, bw, Inches(0.45), ACCENT2)
    add_text(s, x, y + Inches(0.05), bw, Inches(0.4),
             t, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y + Inches(0.55), bw - Inches(0.3),
             bh - Inches(0.6), body, size=12, color=INK,
             align=PP_ALIGN.CENTER)
    if i < 3:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   x + bw + Inches(0.01),
                                   y + bh / 2 - Inches(0.15),
                                   Inches(0.18), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

add_bullets(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.2),
            ["Best for: documentation Q&A, internal knowledge, fresh data",
             "Always cite sources — \"Answer using only the context. Cite the source ID.\"",
             "Quality of retrieval >> cleverness of the prompt",
             "Watch context length: rerank and truncate aggressively",
             "Evaluate end-to-end: answer accuracy, faithfulness, citation correctness"],
            size=13)


# =====================================================================
# 17. Context engineering
# =====================================================================
s = base_slide("Context engineering", eyebrow="Advanced Patterns")

add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.7),
         "Treat the context window as a designed system: every token should earn "
         "its place.",
         size=16, color=INK)

card(s, Inches(0.6), Inches(2.7), Inches(3.95), Inches(3.9),
     "What goes in",
     ["System prompt with role, rules, format",
      "Recent conversation, summarized older turns",
      "Tool definitions and schemas",
      "Retrieved documents with IDs",
      "Few-shot examples when needed",
      "Current user input last"])
card(s, Inches(4.7), Inches(2.7), Inches(3.95), Inches(3.9),
     "Discipline",
     ["Short, ordered, deduplicated",
      "Trim before you grow",
      "Pin format with delimiters",
      "Prefer instructions over hints",
      "Keep a stable structure across calls",
      "Budget tokens per section"])
card(s, Inches(8.8), Inches(2.7), Inches(3.95), Inches(3.9),
     "Failure modes",
     ["Lost-in-the-middle attention",
      "Conflicting instructions",
      "Stale conversation history",
      "Leaking irrelevant retrieval hits",
      "Overstuffed system prompt",
      "Implicit assumptions never stated"])


# =====================================================================
# 18. Common pitfalls
# =====================================================================
s = base_slide("Common pitfalls", eyebrow="Patterns to Avoid")

pitfalls = [
    ("Vague instructions",
     "\"Make it better\" → no measurable target. Specify what better means."),
    ("Negative-only constraints",
     "Telling the model what not to do without giving a fallback."),
    ("Overstuffed prompts",
     "Long, contradictory context dilutes attention and increases cost."),
    ("Mixing roles",
     "System persona conflicts with task instructions or examples."),
    ("Untested edge cases",
     "Empty input, hostile input, multilingual input, very long input."),
    ("No format contract",
     "Free-form output breaks downstream parsers. Pin a schema."),
]
for i, (t, body) in enumerate(pitfalls):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.15)
    y = Inches(1.85 + row * 1.7)
    add_rect(s, x, y, Inches(5.95), Inches(1.55), WHITE, line=LINE)
    # Red dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25),
                             y + Inches(0.25), Inches(0.28), Inches(0.28))
    dot.fill.solid()
    dot.fill.fore_color.rgb = RGBColor(0xE0, 0x3A, 0x3A)
    dot.line.fill.background()
    add_text(s, x + Inches(0.65), y + Inches(0.15), Inches(5.2), Inches(0.45),
             t, size=15, bold=True, color=NAVY)
    add_text(s, x + Inches(0.65), y + Inches(0.65), Inches(5.2), Inches(0.85),
             body, size=12, color=INK)


# =====================================================================
# 19. Risks: hallucination, bias, injection
# =====================================================================
s = base_slide("Risks to design against", eyebrow="Risks & Safety")

risks = [
    ("Hallucination",
     "The model invents facts that sound right.",
     "Ground in retrieval. Require citations. Ask for 'I don't know.'"),
    ("Bias",
     "Outputs reflect training-data skew on sensitive topics.",
     "Audit on representative data. Constrain scope. Add review."),
    ("Prompt injection",
     "Untrusted text in the context overrides your instructions.",
     "Treat all inputs as untrusted. Separate trust zones in the prompt."),
    ("Prompt leaking",
     "Model reveals your system prompt or proprietary instructions.",
     "Don't put secrets in prompts. Add policy + monitoring."),
    ("Jailbreaking",
     "Crafted inputs bypass safety guardrails.",
     "Layered defenses: model, prompt, application, post-filter."),
    ("Data leakage",
     "PII or confidential content escapes via outputs or logs.",
     "Redact at ingest. Minimize retention. Encrypt and access-control."),
]
for i, (t, body, mitig) in enumerate(risks):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.15)
    y = Inches(1.85 + row * 2.6)
    add_rect(s, x, y, Inches(3.95), Inches(2.4), WHITE, line=LINE)
    add_rect(s, x, y, Inches(3.95), Inches(0.45), RGBColor(0xE0, 0x3A, 0x3A))
    add_text(s, x + Inches(0.2), y + Inches(0.05), Inches(3.6), Inches(0.4),
             t, size=14, bold=True, color=WHITE)
    add_text(s, x + Inches(0.2), y + Inches(0.55), Inches(3.6), Inches(0.7),
             body, size=12, color=INK)
    add_text(s, x + Inches(0.2), y + Inches(1.35), Inches(3.6), Inches(0.95),
             "→ " + mitig, size=11, bold=True, color=ACCENT2)


# =====================================================================
# 20. Defense strategies
# =====================================================================
s = base_slide("Defense strategies", eyebrow="Risks & Safety")

defenses = [
    ("Trust zones",
     "Mark which parts of the context are user-controlled vs. system-controlled. "
     "Never let user text be parsed as instructions."),
    ("Input/output filters",
     "Validate inputs against a schema. Validate outputs against a parser. "
     "Reject or repair on failure."),
    ("Allowlists & schemas",
     "Constrain answers to known values, JSON shapes, or function signatures "
     "via tool-use APIs."),
    ("Defense in depth",
     "Combine model-side guardrails with application-level checks and "
     "post-processing review."),
    ("Red-team continuously",
     "Maintain an evolving suite of injection and jailbreak tests. "
     "Replay them on every prompt change."),
    ("Human in the loop",
     "Require approval for high-impact actions. "
     "Confidence-based routing for ambiguous cases."),
]
for i, (t, body) in enumerate(defenses):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.15)
    y = Inches(1.85 + row * 1.7)
    add_rect(s, x, y, Inches(5.95), Inches(1.55), WHITE, line=LINE)
    # Green check
    check = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25),
                               y + Inches(0.25), Inches(0.28), Inches(0.28))
    check.fill.solid()
    check.fill.fore_color.rgb = ACCENT2
    check.line.fill.background()
    add_text(s, x + Inches(0.65), y + Inches(0.15), Inches(5.2), Inches(0.45),
             t, size=15, bold=True, color=NAVY)
    add_text(s, x + Inches(0.65), y + Inches(0.65), Inches(5.2), Inches(0.85),
             body, size=12, color=INK)


# =====================================================================
# 21. Iterative workflow
# =====================================================================
s = base_slide("An iterative prompt workflow", eyebrow="Workflow & Evaluation")

steps = [
    ("Define",   "Goal, inputs, outputs,\nsuccess metric"),
    ("Draft",    "Minimal prompt; pin\nsettings; no examples"),
    ("Test",     "Run on a small,\nrepresentative set"),
    ("Diagnose", "Cluster failures. What\nkind of error is it?"),
    ("Improve",  "Add the smallest fix:\ninstruction, example, tool"),
    ("Lock",     "Versioned prompt,\neval suite, rollback"),
]
y = Inches(2.0); bw = Inches(1.95); bh = Inches(1.5); gap = Inches(0.13)
for i, (t, body) in enumerate(steps):
    x = Inches(0.6) + (bw + gap) * i
    add_rect(s, x, y, bw, bh, WHITE, line=LINE)
    add_rect(s, x, y, bw, Inches(0.4), NAVY)
    add_text(s, x, y + Inches(0.04), bw, Inches(0.35),
             t, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(0.5), bw - Inches(0.2),
             bh - Inches(0.55), body, size=11, color=INK,
             align=PP_ALIGN.CENTER)
    if i < 5:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   x + bw + Inches(0.005),
                                   y + bh / 2 - Inches(0.13),
                                   Inches(0.12), Inches(0.26))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

add_text(s, Inches(0.6), Inches(3.9), Inches(12.1), Inches(0.5),
         "Evaluation is the discipline that turns prompting into engineering",
         size=16, bold=True, color=NAVY)

add_bullets(s, Inches(0.6), Inches(4.45), Inches(12.1), Inches(2.5),
            ["Build a small golden set early — 20-50 examples beat zero",
             "Track: accuracy, format compliance, latency, cost, refusal rate",
             "Use LLM-as-judge for subjective metrics, but validate the judge",
             "Treat prompts like code: review, version, and test on every change",
             "Roll out behind feature flags; compare against the previous version"],
            size=13)


# =====================================================================
# 22. Cheat sheet
# =====================================================================
s = base_slide("One-page cheat sheet", eyebrow="Takeaways")

left = [
    "Lead with a strong verb: classify, extract, summarize",
    "State the audience and the format",
    "Show 2-5 diverse examples when format matters",
    "Say what to do, not what to avoid",
    "Use delimiters: ###, XML tags, fenced blocks",
    "Decompose hard tasks into chained prompts",
]
right = [
    "Pin temperature/top-p when measuring changes",
    "Add 'Let's think step by step.' for reasoning",
    "Ground knowledge tasks with retrieval + citations",
    "Treat all input text as untrusted",
    "Validate output against a schema",
    "Version prompts, run an eval suite",
]

add_text(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(0.5),
         "Do this", size=18, bold=True, color=ACCENT2)
add_bullets(s, Inches(0.6), Inches(2.4), Inches(6.0), Inches(4.5),
            left, size=14)

add_text(s, Inches(6.8), Inches(1.8), Inches(6.0), Inches(0.5),
         "And this", size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(6.8), Inches(2.4), Inches(6.0), Inches(4.5),
            right, size=14)


# =====================================================================
# 23. Resources
# =====================================================================
s = base_slide("Resources & next steps", eyebrow="Keep learning")

resources = [
    ("Prompt Engineering Guide",
     "promptingguide.ai",
     "Comprehensive reference: techniques, models, papers, prompt hub."),
    ("OpenAI / Anthropic / AWS docs",
     "Provider best-practice guides",
     "Model-specific patterns, system prompt advice, tool-use examples."),
    ("Hands-on workshop",
     "This repo: Streamlit + Bedrock",
     "Foundations, zero/few-shot, CoT, advanced, agentic, safety, RAG."),
    ("Papers worth reading",
     "CoT, Self-Consistency, ReAct, ToT, RAG",
     "Original sources for the techniques covered in this session."),
]
for i, (t, where, body) in enumerate(resources):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.15)
    y = Inches(1.85 + row * 2.4)
    add_rect(s, x, y, Inches(5.95), Inches(2.2), WHITE, line=LINE)
    add_rect(s, x, y, Inches(0.08), Inches(2.2), ACCENT)
    add_text(s, x + Inches(0.25), y + Inches(0.15), Inches(5.6), Inches(0.45),
             t, size=16, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.6), Inches(5.6), Inches(0.4),
             where, size=12, bold=True, color=ACCENT2)
    add_text(s, x + Inches(0.25), y + Inches(1.05), Inches(5.6), Inches(1.0),
             body, size=12, color=INK)


# =====================================================================
# 24. Q&A / closing
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
add_rect(s, 0, Inches(7.0), prs.slide_width, Inches(0.5), ACCENT)
add_rect(s, 0, 0, Inches(0.4), prs.slide_height, ACCENT2)

add_text(s, Inches(1.0), Inches(2.4), Inches(11), Inches(1.2),
         "Questions?", size=72, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(3.7), Inches(11), Inches(0.7),
         "Let's discuss the patterns that fit your use case.",
         size=22, color=RGBColor(0xCB, 0xD5, 0xE1))

add_text(s, Inches(1.0), Inches(5.5), Inches(11), Inches(0.5),
         "Thanks for joining the workshop", size=18, bold=True, color=ACCENT)
add_text(s, Inches(1.0), Inches(6.0), Inches(11), Inches(0.4),
         "Reference: promptingguide.ai", size=14, color=WHITE)


# Save
out = "/Users/erictole/demo/prompt-engineering-workshop/prompt_engineering_workshop.pptx"
prs.save(out)
print(f"Wrote {out}  ·  {len(prs.slides)} slides")
